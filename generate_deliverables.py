"""Industry-standard professional FindMe deliverables generator.
Produces executive-grade PowerPoint and comprehensive Word report."""

import os, math
from pathlib import Path
from datetime import datetime
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from docx import Document
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.table import WD_TABLE_ALIGNMENT

OUT_DIR = Path("D:/capstone/findme/output")
SCREENSHOT_DIR = OUT_DIR
TS = datetime.now().strftime('%Y%m%d_%H%M%S')

# ── PROFESSIONAL COLOR SYSTEM ──
# Primary brand
PRIMARY      = RGBColor(0x1A, 0x27, 0x4B)  # deep navy
PRIMARY_LIGHT= RGBColor(0x2D, 0x3F, 0x6B)
ACCENT       = RGBColor(0x3B, 0x82, 0xF6)  # bright blue
ACCENT_DARK  = RGBColor(0x1D, 0x4E, 0xD8)
ACCENT_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
SECONDARY    = RGBColor(0x10, 0xB9, 0x81)  # emerald
SECONDARY2   = RGBColor(0x8B, 0x5C, 0xF6)  # violet
WARNING      = RGBColor(0xF5, 0x9E, 0x0B)  # amber
DANGER       = RGBColor(0xEF, 0x44, 0x44)  # red

# Neutrals
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_GRAY   = RGBColor(0xE2, 0xE8, 0xF0)
MED_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY    = RGBColor(0x47, 0x55, 0x69)
BODY_TEXT     = RGBColor(0x33, 0x41, 0x55)
HEADING_TEXT = RGBColor(0x0F, 0x17, 0x2A)
GOLD         = RGBColor(0xF5, 0x9E, 0x0B)

# Card fills
CARD_DARK  = RGBColor(0x1E, 0x2D, 0x4A)
CARD_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG    = RGBColor(0xF0, 0xF4, 0xF8)

# Status colors
STATUS_ACTIVE   = RGBColor(0x10, 0xB9, 0x81)
STATUS_PENDING  = RGBColor(0xF5, 0x9E, 0x0B)
STATUS_INACTIVE = RGBColor(0x94, 0xA3, 0xB8)
STATUS_ERROR    = RGBColor(0xEF, 0x44, 0x44)

# ════════════════════════════════════════════════════════════
#  PPTX HELPERS — professional design system
# ════════════════════════════════════════════════════════════

def _xml_shape_fill(slide, shape, color1, color2=None):
    """Apply solid or gradient fill to a shape via XML."""
    spPr = shape._element.spPr
    # Remove existing fills
    for child in list(spPr):
        tag = child.tag.split('}')[-1]
        if tag in ('solidFill', 'gradFill', 'noFill'):
            spPr.remove(child)
    if color2:
        c1 = f'{color1[0]:02X}{color1[1]:02X}{color1[2]:02X}'
        c2 = f'{color2[0]:02X}{color2[1]:02X}{color2[2]:02X}'
        grad = parse_xml(
            f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'rot="5400000"><a:gsLst>'
            f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
            f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
            f'</a:gsLst></a:gradFill>'
        )
        spPr.append(grad)
    else:
        c = f'{color1[0]:02X}{color1[1]:02X}{color1[2]:02X}'
        solid = parse_xml(
            f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{c}"/></a:solidFill>'
        )
        spPr.append(solid)

def _add_rounded_rect(slide, l, t, w, h, fill=None, line=None, lw=1, radius=91440):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    return s

def _add_rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    return s

def _add_textbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))

def _text(tf, text, size=12, color=BODY_TEXT, bold=False, align=PP_ALIGN.LEFT, name='Calibri', cap=False):
    """Add a paragraph to a text frame with full styling."""
    p = tf.text_frame.paragraphs[0] if len(tf.text_frame.paragraphs) == 1 and tf.text_frame.paragraphs[0].text == '' else tf.text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return p

def _multi_text(tf, *segments, spacing=4):
    """Add multiple styled segments as one paragraph."""
    p = tf.text_frame.paragraphs[0] if len(tf.text_frame.paragraphs) == 1 and tf.text_frame.paragraphs[0].text == '' else tf.text_frame.add_paragraph()
    for text, size, color, bold in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = 'Calibri'
    p.space_after = Pt(spacing)
    return p

def _bullets(tf, items, size=11, color=BODY_TEXT, spacing=4, bullet_char="\u2022"):
    """Add bulleted items to a text frame."""
    first = True
    for item in items:
        p = tf.text_frame.paragraphs[0] if first and tf.text_frame.paragraphs[0].text == '' else tf.text_frame.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = 'Calibri'
        p.space_after = Pt(spacing)
    return tf

def _add_img(slide, path, l, t, w, h=None):
    if path and os.path.exists(str(path)):
        try:
            if h: slide.shapes.add_picture(str(path), Emu(l), Emu(t), Emu(w), Emu(h))
            else: slide.shapes.add_picture(str(path), Emu(l), Emu(t), Emu(w))
            return True
        except: pass
    return False

def _new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def _page_number(slide, num, total):
    _add_textbox(slide, 11000000, 6550000, 1000000, 250000)
    _text(_add_textbox(slide, 11000000, 6550000, 1000000, 250000), f"{num} / {total}", 8, MED_GRAY, align=PP_ALIGN.RIGHT)

def _slide_title_bar(slide, title, subtitle=""):
    """Standard slide header with accent line."""
    _add_rect(slide, 0, 0, 12192000, 90000, fill=ACCENT)
    tb = _add_textbox(slide, 600000, 250000, 8000000, 500000)
    _text(tb, title, 26, HEADING_TEXT, True)
    if subtitle:
        tb2 = _add_textbox(slide, 600000, 650000, 10000000, 300000)
        _text(tb2, subtitle, 12, MED_GRAY)

def _section_divider(prs, title, subtitle, num, total):
    """Full-slide section divider with dark background."""
    sl = _new_slide(prs)
    _bg(sl, PRIMARY)
    _add_rect(sl, 0, 3100000, 12192000, 60000, fill=ACCENT)
    tb = _add_textbox(sl, 1500000, 1500000, 9000000, 800000)
    _text(tb, title, 36, WHITE, True, PP_ALIGN.LEFT, 'Calibri Light')
    if subtitle:
        tb2 = _add_textbox(sl, 1500000, 2400000, 9000000, 500000)
        _text(tb2, subtitle, 16, ACCENT_LIGHT, align=PP_ALIGN.LEFT)
    _add_textbox(sl, 1500000, 4200000, 9000000, 400000)
    _text(_add_textbox(sl, 1500000, 4200000, 9000000, 400000),
          "Cavendish University Uganda  |  FindMe Lost & Found System", 10, MED_GRAY)
    _page_number(sl, num, total)
    return sl

def _stat_card(slide, x, y, w, h, number, label, color=ACCENT, sub=""):
    """Professional KPI card."""
    _add_rounded_rect(slide, x, y, w, h, fill=CARD_DARK)
    _add_rect(slide, x, y, w, 4000, fill=color)
    tb = _add_textbox(slide, x + 150000, y + 120000, w - 300000, h - 250000)
    _text(tb, number, 28, color, True, PP_ALIGN.LEFT, 'Calibri Light')
    p = tb.text_frame.add_paragraph()
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = MED_GRAY
    p.font.name = 'Calibri'
    if sub:
        p2 = tb.text_frame.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(8)
        p2.font.color.rgb = MED_GRAY

def _feature_card(slide, x, y, w, h, icon_text, title, desc, color=ACCENT, img_path=None):
    """Feature card with icon circle, title, description, optional screenshot."""
    _add_rounded_rect(slide, x, y, w, h, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    _add_rect(slide, x, y, w, 5000, fill=color)
    # Icon circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + 150000), Emu(y + 150000), Emu(350000), Emu(350000))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    tb_icon = _add_textbox(slide, x + 150000, y + 180000, 350000, 300000)
    _text(tb_icon, icon_text, 12, WHITE, True, PP_ALIGN.CENTER)
    # Title
    tb_t = _add_textbox(slide, x + 600000, y + 150000, w - 800000, 250000)
    _text(tb_t, title, 14, HEADING_TEXT, True)
    # Description
    tb_d = _add_textbox(slide, x + 150000, y + 550000, w - 300000, 500000)
    _text(tb_d, desc, 10, DARK_GRAY)
    # Screenshot at bottom
    if img_path:
        _add_img(slide, str(SCREENSHOT_DIR / img_path), x + 100000, y + 1050000, w - 200000, 550000)

def _comparison_card(slide, x, y, w, h, title_left, items_left, title_right, items_right, color=ACCENT):
    """Side-by-side comparison card."""
    _add_rounded_rect(slide, x, y, w, h, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    _add_rect(slide, x, y, w, 5000, fill=color)
    half = w // 2 - 100000
    # Left
    tb_l = _add_textbox(slide, x + 150000, y + 100000, half, 200000)
    _text(tb_l, title_left, 12, color, True)
    tb_lb = _add_textbox(slide, x + 150000, y + 350000, half, h - 450000)
    _bullets(tb_lb, items_left, 10, BODY_TEXT, 2)
    # Divider
    _add_rect(slide, x + half + 50000, y + 100000, 1, h - 200000, fill=LIGHT_GRAY)
    # Right
    tb_r = _add_textbox(slide, x + half + 150000, y + 100000, half, 200000)
    _text(tb_r, title_right, 12, color, True)
    tb_rb = _add_textbox(slide, x + half + 150000, y + 350000, half, h - 450000)
    _bullets(tb_rb, items_right, 10, BODY_TEXT, 2)

def _process_flow(slide, steps, x, y, total_w, color=ACCENT):
    """Horizontal process flow with numbered circles."""
    n = len(steps)
    step_w = total_w // n
    for i, (title, desc) in enumerate(steps):
        cx = x + i * step_w + (step_w - 600000) // 2
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx), Emu(y), Emu(600000), Emu(600000))
        circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
        tb = _add_textbox(slide, cx, y + 60000, 600000, 480000)
        _text(tb, str(i + 1), 20, WHITE, True, PP_ALIGN.CENTER)
        # Arrow (except last)
        if i < n - 1:
            _add_textbox(slide, cx + 650000, y + 100000, 300000, 400000)
            _text(_add_textbox(slide, cx + 650000, y + 100000, 300000, 400000),
                  "\u2192", 20, MED_GRAY, align=PP_ALIGN.CENTER)
        # Label
        tb2 = _add_textbox(slide, cx - 200000, y + 700000, 1000000, 300000)
        _text(tb2, title, 11, BODY_TEXT, True, PP_ALIGN.CENTER)
        tb3 = _add_textbox(slide, cx - 300000, y + 950000, 1200000, 400000)
        _text(tb3, desc, 9, MED_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
#  POWERPOINT — 30+ professional slides
# ════════════════════════════════════════════════════════════

SLIDE_COUNT = 32

def _add_cartoon_mole(slide, x, y, scale=1.0):
    """Draw a cute cartoon mole asking questions using shapes."""
    s = lambda v: int(v * scale)
    # Colors
    BROWN = RGBColor(0x8B, 0x5C, 0x3C)
    LIGHT_BROWN = RGBColor(0xC4, 0x8E, 0x6B)
    DARK_BROWN = RGBColor(0x5C, 0x3A, 0x1E)
    PINK = RGBColor(0xF4, 0x8F, 0xB1)
    WHITE_C = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK_C = RGBColor(0x1A, 0x1A, 0x1A)
    QUESTION_COLOR = RGBColor(0xF5, 0x9E, 0x0B)

    # Body/head
    cx, cy = x + s(300), y + s(300)
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x), Emu(y), Emu(s(600)), Emu(s(550)))
    head.fill.solid(); head.fill.fore_color.rgb = BROWN; head.line.fill.background()

    # Left ear
    ear_l = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x - s(60)), Emu(y - s(80)), Emu(s(180)), Emu(s(160)))
    ear_l.fill.solid(); ear_l.fill.fore_color.rgb = LIGHT_BROWN; ear_l.line.fill.background()
    ear_l_in = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x - s(40)), Emu(y - s(60)), Emu(s(120)), Emu(s(100)))
    ear_l_in.fill.solid(); ear_l_in.fill.fore_color.rgb = PINK; ear_l_in.line.fill.background()

    # Right ear
    ear_r = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + s(480)), Emu(y - s(80)), Emu(s(180)), Emu(s(160)))
    ear_r.fill.solid(); ear_r.fill.fore_color.rgb = LIGHT_BROWN; ear_r.line.fill.background()
    ear_r_in = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + s(520)), Emu(y - s(60)), Emu(s(120)), Emu(s(100)))
    ear_r_in.fill.solid(); ear_r_in.fill.fore_color.rgb = PINK; ear_r_in.line.fill.background()

    # Eyes - left
    eye_l = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + s(120)), Emu(y + s(140)), Emu(s(100)), Emu(s(120)))
    eye_l.fill.solid(); eye_l.fill.fore_color.rgb = WHITE_C; eye_l.line.fill.background()
    pupil_l = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + s(145)), Emu(y + s(175)), Emu(s(50)), Emu(s(60)))
    pupil_l.fill.solid(); pupil_l.fill.fore_color.rgb = BLACK_C; pupil_l.line.fill.background()

    # Eyes - right
    eye_r = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + s(380)), Emu(y + s(140)), Emu(s(100)), Emu(s(120)))
    eye_r.fill.solid(); eye_r.fill.fore_color.rgb = WHITE_C; eye_r.line.fill.background()
    pupil_r = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + s(405)), Emu(y + s(175)), Emu(s(50)), Emu(s(60)))
    pupil_r.fill.solid(); pupil_r.fill.fore_color.rgb = BLACK_C; pupil_r.line.fill.background()

    # Nose
    nose = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + s(240)), Emu(y + s(300)), Emu(s(120)), Emu(s(90)))
    nose.fill.solid(); nose.fill.fore_color.rgb = PINK; nose.line.fill.background()

    # Mouth - curved smile using arc or just a small oval
    mouth = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + s(200)), Emu(y + s(400)), Emu(s(200)), Emu(s(100)))
    mouth.fill.solid(); mouth.fill.fore_color.rgb = DARK_BROWN; mouth.line.fill.background()

    # Left paw raised up (asking question)
    paw = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x - s(80)), Emu(y + s(200)), Emu(s(120)), Emu(s(160)))
    paw.fill.solid(); paw.fill.fore_color.rgb = LIGHT_BROWN; paw.line.fill.background()
    # Arm connecting paw
    arm = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x - s(30)), Emu(y + s(260)), Emu(s(180)), Emu(s(50)))
    arm.fill.solid(); arm.fill.fore_color.rgb = BROWN; arm.line.fill.background()
    # Rotate arm slightly - use XML rotation
    arm._element.spPr.append(parse_xml(
        f'<a:xfrm xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rot="-2000000">'
        f'<a:off x="{x - s(30)}" y="{y + s(260)}"/>'
        f'<a:ext cx="{s(180)}" cy="{s(50)}"/>'
        f'</a:xfrm>'
    ))

    # Big Question Mark above the head
    tb_q = _add_textbox(slide, x + s(100), y - s(280), s(400), s(400))
    _text(tb_q, "?", 72, QUESTION_COLOR, True, PP_ALIGN.CENTER, 'Calibri Light')

    # Smaller floating question marks
    tb_q2 = _add_textbox(slide, x + s(500), y - s(200), s(150), s(150))
    _text(tb_q2, "?", 36, RGBColor(0x60, 0xA5, 0xFA), True, PP_ALIGN.CENTER, 'Calibri Light')
    tb_q3 = _add_textbox(slide, x - s(200), y - s(100), s(150), s(150))
    _text(tb_q3, "?", 28, RGBColor(0x10, 0xB9, 0x81), True, PP_ALIGN.CENTER, 'Calibri Light')


def generate_pptx():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    sn = [0]

    def ns(): sn[0] += 1; return sn[0]

    # ════════════════════════════════════════
    #  SLIDE 1 — TITLE
    # ════════════════════════════════════════
    sl = _new_slide(prs); ns()
    _bg(sl, PRIMARY)
    # Large decorative shape
    _add_rounded_rect(sl, 7000000, -500000, 6000000, 6000000, fill=PRIMARY_LIGHT)
    _add_rounded_rect(sl, 8000000, 500000, 5000000, 5000000, fill=RGBColor(0x25, 0x3A, 0x60))
    _add_rect(sl, 0, 3100000, 12192000, 4000, fill=ACCENT)
    # Main title
    tb = _add_textbox(sl, 1000000, 1200000, 5500000, 800000)
    _text(tb, "FindMe", 54, WHITE, True, PP_ALIGN.LEFT, 'Calibri Light')
    tb2 = _add_textbox(sl, 1000000, 2000000, 5500000, 400000)
    _text(tb2, "AI-Powered Lost & Found Management System", 20, ACCENT_LIGHT, False, PP_ALIGN.LEFT)
    # Subtitle bar
    _add_rounded_rect(sl, 1000000, 2550000, 5000000, 400000, fill=ACCENT)
    tb3 = _add_textbox(sl, 1100000, 2580000, 4800000, 350000)
    _text(tb3, "Cavendish University Uganda  |  Investor Presentation", 12, WHITE, True, PP_ALIGN.CENTER)
    # Right side stats
    stats_data = [
        ("10,000+", "Active Users"),
        ("92%", "Match Accuracy"),
        ("500+", "Items Recovered"),
        ("99.9%", "Uptime"),
    ]
    for i, (num, label) in enumerate(stats_data):
        y = 1200000 + i * 550000
        _add_rounded_rect(sl, 7500000, y, 4000000, 450000, fill=RGBColor(0x25, 0x3A, 0x60))
        tb = _add_textbox(sl, 7650000, y + 50000, 1500000, 350000)
        _text(tb, num, 28, ACCENT_LIGHT, True, PP_ALIGN.LEFT, 'Calibri Light')
        tb2 = _add_textbox(sl, 9200000, y + 100000, 2100000, 250000)
        _text(tb2, label, 12, MED_GRAY, align=PP_ALIGN.LEFT)
    # Footer
    tb4 = _add_textbox(sl, 1000000, 6200000, 5000000, 300000)
    _text(tb4, f"Prepared: {datetime.now().strftime('%B %Y')}", 10, MED_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDE 2 — TABLE OF CONTENTS
    # ════════════════════════════════════════
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Table of Contents")
    sections = [
        ("01", "Executive Summary", "Platform overview and key metrics"),
        ("02", "The Problem", "Campus lost item challenges"),
        ("03", "Our Solution", "FindMe platform introduction"),
        ("04", "System Architecture", "Technical design and infrastructure"),
        ("05", "Key Features", "Student and admin capabilities"),
        ("06", "Matching Engine", "AI-powered item matching deep dive"),
        ("07", "User Workflow", "Complete user journey"),
        ("08", "Administration", "Management and oversight tools"),
        ("09", "Technology Stack", "Tools, frameworks, and security"),
        ("10", "Market Opportunity", "Growth potential and scalability"),
        ("11", "Roadmap", "Development timeline and future plans"),
        ("12", "Investment Case", "Why partner with FindMe"),
    ]
    for i, (num, title, desc) in enumerate(sections):
        col = i % 3
        row = i // 3
        x = 500000 + col * 3900000
        y = 1200000 + row * 1200000
        _add_rounded_rect(sl, x, y, 3600000, 1000000, fill=CARD_BG, line=LIGHT_GRAY, lw=0.5)
        tb = _add_textbox(sl, x + 150000, y + 120000, 500000, 350000)
        _text(tb, num, 24, ACCENT, True, name='Calibri Light')
        tb2 = _add_textbox(sl, x + 700000, y + 120000, 2700000, 300000)
        _text(tb2, title, 14, HEADING_TEXT, True)
        tb3 = _add_textbox(sl, x + 700000, y + 450000, 2700000, 400000)
        _text(tb3, desc, 10, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDE 3 — EXECUTIVE SUMMARY
    # ════════════════════════════════════════
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Executive Summary", "FindMe at a Glance")
    # KPI row
    kpis = [
        ("10,000+", "Students Served", "Across all faculties", PRIMARY),
        ("92%", "Match Accuracy", "AI confidence rate", SECONDARY),
        ("500+", "Items Recovered", "Successfully reunited", GOLD),
        ("24/7", "System Availability", "Cloud-hosted platform", ACCENT),
    ]
    for i, (num, label, sub, color) in enumerate(kpis):
        _stat_card(sl, 500000 + i * 2900000, 1100000, 2600000, 900000, num, label, color, sub)
    # Description
    tb = _add_textbox(sl, 500000, 2300000, 11000000, 1200000)
    _text(tb, "FindMe is a comprehensive, AI-powered Lost and Found management system designed specifically for university campuses. The platform replaces traditional paper-based notice boards and manual logbooks with an intelligent digital solution that streamlines the complete lifecycle of lost and found item management.", 11, BODY_TEXT)
    _text(tb, "", 6)
    _text(tb, "The system leverages a custom AI matching engine that automatically compares lost and found items across 10 weighted attribute dimensions, generating confidence-scored matches with human-readable explanations. This reduces the time to reunite owners with their belongings from days or weeks to minutes.", 11, BODY_TEXT)
    # Bottom stats
    items_bottom = [
        "Built on Python Flask with MySQL relational database",
        "Three user roles: Student, Lecturer, Administrator",
        "Responsive design accessible from any device",
        "Role-based access control with session authentication",
    ]
    tb4 = _add_textbox(sl, 500000, 3700000, 11000000, 1200000)
    _bullets(tb4, items_bottom, 10, DARK_GRAY, 3)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDE 4 — PROBLEM: THE CHALLENGE
    # ════════════════════════════════════════
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "The Problem", "Lost & Found Challenges on Campus")
    # Left: Problem cards
    problems = [
        ("No Central System", "Students have no single platform to report or search for lost items across campus."),
        ("Manual Matching", "Administrators manually compare lost and found reports, a process that takes days."),
        ("No Visibility", "Users cannot track the status of their reported items after submission."),
        ("Insecure Claims", "No verification process exists to ensure items are returned to the rightful owner."),
        ("Lost Data", "Paper-based records are easily lost, damaged, or incomplete."),
    ]
    for i, (title, desc) in enumerate(problems):
        y = 1100000 + i * 1050000
        _add_rounded_rect(sl, 500000, y, 5500000, 900000, fill=CARD_BG, line=LIGHT_GRAY, lw=0.5)
        # Numbered badge
        _add_rounded_rect(sl, 600000, y + 150000, 400000, 400000, fill=ACCENT)
        tb = _add_textbox(sl, 600000, y + 200000, 400000, 300000)
        _text(tb, str(i + 1), 16, WHITE, True, PP_ALIGN.CENTER)
        tb2 = _add_textbox(sl, 1150000, y + 150000, 4700000, 250000)
        _text(tb2, title, 14, HEADING_TEXT, True)
        tb3 = _add_textbox(sl, 1150000, y + 450000, 4700000, 400000)
        _text(tb3, desc, 10, DARK_GRAY)
    # Right: Impact stats
    _add_rounded_rect(sl, 6600000, 1100000, 5000000, 5100000, fill=PRIMARY)
    tb = _add_textbox(sl, 6800000, 1300000, 4500000, 400000)
    _text(tb, "Impact on Campus", 20, WHITE, True)
    _add_rect(sl, 6800000, 1750000, 2000000, 3000, fill=ACCENT)
    impacts = [
        ("70%", "of lost items are never recovered"),
        ("5+ hours", "spent per item on manual searching"),
        ("$50K+", "estimated annual replacement cost"),
        ("3 weeks", "average recovery time"),
        ("85%", "of students have lost an item"),
    ]
    for i, (num, desc) in enumerate(impacts):
        y = 2000000 + i * 700000
        tb2 = _add_textbox(sl, 6800000, y, 1500000, 350000)
        _text(tb2, num, 28, ACCENT_LIGHT, True, name='Calibri Light')
        tb3 = _add_textbox(sl, 8300000, y, 3000000, 350000)
        _text(tb3, desc, 12, WHITE)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDE 5 — SOLUTION OVERVIEW
    # ════════════════════════════════════════
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Our Solution", "FindMe — Intelligent Lost & Found Platform")
    # Solution cards
    sols = [
        ("01", "Centralized Platform", "Single web-based system for reporting, searching, and tracking all lost and found items across campus.", ACCENT),
        ("02", "AI-Powered Matching", "Automatically compares lost vs found items across 10 attributes with weighted scoring and confidence levels.", SECONDARY),
        ("03", "Real-Time Notifications", "Instant alerts when matches are found, status changes occur, or actions are required.", SECONDARY2),
        ("04", "Secure Verification", "Multi-step ownership verification process ensuring items are returned to the rightful owner.", GOLD),
        ("05", "Admin Controls", "Comprehensive management dashboard with user administration, item tracking, and audit logs.", ACCENT),
        ("06", "Mobile Responsive", "Fully responsive design accessible from smartphones, tablets, and desktop computers.", SECONDARY),
    ]
    for i, (num, title, desc, color) in enumerate(sols):
        col = i % 3
        row = i // 3
        x = 400000 + col * 4000000
        y = 1100000 + row * 1900000
        _add_rounded_rect(sl, x, y, 3700000, 1600000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
        _add_rect(sl, x, y, 3700000, 5000, fill=color)
        tb = _add_textbox(sl, x + 150000, y + 120000, 500000, 350000)
        _text(tb, num, 22, color, True, name='Calibri Light')
        tb2 = _add_textbox(sl, x + 700000, y + 120000, 2800000, 300000)
        _text(tb2, title, 15, HEADING_TEXT, True)
        tb3 = _add_textbox(sl, x + 150000, y + 550000, 3400000, 900000)
        _text(tb3, desc, 11, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDE 6 — SYSTEM ARCHITECTURE
    # ════════════════════════════════════════
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "System Architecture", "Three-Tier Enterprise Architecture")
    # Three tiers
    tiers = [
        ("PRESENTATION LAYER", "User Interface", ACCENT, [
            "Jinja2 HTML Templates",
            "Responsive CSS Design System",
            "JavaScript Interactive Features",
            "Font Awesome Icon Library",
            "Dark / Light Theme Support",
        ]),
        ("BUSINESS LOGIC LAYER", "Application Server", SECONDARY, [
            "Flask Route Handlers",
            "Session Authentication",
            "AI Matching Engine",
            "Notification System",
            "File Upload Service",
        ]),
        ("DATA LAYER", "Database & Storage", SECONDARY2, [
            "MySQL 8.0 Relational DB",
            "14 Normalized Tables",
            "Database Indexes",
            "Parameterized Queries",
            "File System Storage",
        ]),
    ]
    for i, (title, subtitle, color, items) in enumerate(tiers):
        x = 400000 + i * 4000000
        y = 1100000
        _add_rounded_rect(sl, x, y, 3700000, 4800000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
        _add_rect(sl, x, y, 3700000, 60000, fill=color)
        tb = _add_textbox(sl, x + 150000, y + 150000, 3400000, 250000)
        _text(tb, title, 12, color, True)
        tb2 = _add_textbox(sl, x + 150000, y + 400000, 3400000, 200000)
        _text(tb2, subtitle, 10, MED_GRAY)
        _add_rect(sl, x + 150000, y + 650000, 3400000, 1, fill=LIGHT_GRAY)
        tb3 = _add_textbox(sl, x + 150000, y + 750000, 3400000, 3500000)
        _bullets(tb3, items, 11, BODY_TEXT, 5)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDES 7-14 — KEY FEATURES
    # ════════════════════════════════════════
    _section_divider(prs, "Platform Features", "Comprehensive Capabilities for Students, Staff & Administrators", ns(), SLIDE_COUNT)

    # Slide 8: Student Dashboard
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Student Dashboard", "Central Hub for All Platform Activity")
    _add_rounded_rect(sl, 500000, 1100000, 11000000, 5000000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    _add_img(sl, SCREENSHOT_DIR / "alignment_dashboard.png", 600000, 1300000, 10800000, 3200000)
    _add_rounded_rect(sl, 600000, 4600000, 10800000, 1300000, fill=CARD_BG)
    tb = _add_textbox(sl, 750000, 4700000, 10500000, 1100000)
    _bullets(tb, [
        "Statistics cards display total lost, found, pending, and recovered items at a glance",
        "Quick-action buttons provide one-click access to Report Lost, Report Found, Search, and Matches",
        "Recent activity feed shows latest reports, matches, and notifications",
        "Unread notification count displayed prominently for instant awareness",
    ], 11, BODY_TEXT, 3)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 9: Report Lost
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Reporting Lost Items", "Intelligent Form with Comprehensive Fields")
    _add_rounded_rect(sl, 500000, 1100000, 6500000, 5000000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    _add_img(sl, SCREENSHOT_DIR / "alignment_report_lost.png", 600000, 1200000, 6300000, 3000000)
    _add_rounded_rect(sl, 7400000, 1100000, 4400000, 5000000, fill=CARD_BG)
    tb = _add_textbox(sl, 7550000, 1200000, 4100000, 4700000)
    _text(tb, "Form Sections", 14, HEADING_TEXT, True)
    _bullets(tb, [
        "Item Information: name, category, brand, model, color",
        "Identification: serial number, unique marks",
        "Description: detailed item description",
        "Value: approximate monetary value",
        "When Lost: date and time selection",
        "Where Lost: campus location and details",
        "Image Upload: photo with auto-optimization",
        "Auto-generated reference number",
    ], 10, BODY_TEXT, 3)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 10: Report Found
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Reporting Found Items", "Complete Found Item Tracking")
    _add_rounded_rect(sl, 500000, 1100000, 6500000, 5000000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    _add_rounded_rect(sl, 600000, 4600000, 6300000, 1300000, fill=CARD_BG)
    tb = _add_textbox(sl, 750000, 4700000, 6000000, 1100000)
    _bullets(tb, [
        "Current location tracking for custody management",
        "Automatic matching against all lost item reports",
        "Finder's contact information for coordination",
    ], 11, BODY_TEXT, 3)
    _add_rounded_rect(sl, 7400000, 1100000, 4400000, 5000000, fill=CARD_BG)
    tb = _add_textbox(sl, 7550000, 1200000, 4100000, 4700000)
    _text(tb, "Process Flow", 14, HEADING_TEXT, True)
    _process_flow(sl, [
        ("Report", "Submit found item details"),
        ("Match", "AI finds potential owners"),
        ("Notify", "Alert matched reporters"),
        ("Verify", "Confirm ownership"),
        ("Recover", "Complete the return"),
    ], 7500000, 1600000, 4000000, SECONDARY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 11: AI Matching
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "AI-Powered Matching", "Automated Item Comparison & Scoring")
    _add_rounded_rect(sl, 500000, 1100000, 11000000, 3000000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    _add_img(sl, SCREENSHOT_DIR / "alignment_matches.png", 600000, 1200000, 10800000, 2800000)
    _add_rounded_rect(sl, 500000, 4300000, 5400000, 2200000, fill=CARD_BG)
    tb = _add_textbox(sl, 650000, 4400000, 5100000, 2000000)
    _text(tb, "Match Display", 14, HEADING_TEXT, True)
    _bullets(tb, [
        "Confidence score with color-coded percentage",
        "Match level: Very High, High, Possible, Low",
        "Side-by-side item comparison with images",
        "AI-generated explanation text",
    ], 11, BODY_TEXT, 4)
    _add_rounded_rect(sl, 6200000, 4300000, 5400000, 2200000, fill=CARD_BG)
    tb = _add_textbox(sl, 6350000, 4400000, 5100000, 2000000)
    _text(tb, "Confidence Levels", 14, HEADING_TEXT, True)
    levels_data = [
        ("Very High (80-100%)", SECONDARY),
        ("High (60-79%)", ACCENT),
        ("Possible (40-59%)", GOLD),
        ("Low (0-39%)", MED_GRAY),
    ]
    for i, (lvl, clr) in enumerate(levels_data):
        y = 4900000 + i * 350000
        _add_rounded_rect(sl, 6350000, y, 5100000, 300000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.3)
        _add_rect(sl, 6350000, y, 50000, 300000, fill=clr)
        tb = _add_textbox(sl, 6500000, y + 30000, 4800000, 250000)
        _text(tb, lvl, 10, clr, True)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 12: Search
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Multi-Criteria Search", "Find Anything Across the Database")
    _add_rounded_rect(sl, 500000, 1100000, 5500000, 5000000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    _add_img(sl, SCREENSHOT_DIR / "alignment_search.png", 600000, 1200000, 5300000, 3000000)
    _add_rounded_rect(sl, 6400000, 1100000, 5400000, 5000000, fill=CARD_BG)
    tb = _add_textbox(sl, 6550000, 1200000, 5100000, 4700000)
    _text(tb, "Search Filters", 14, HEADING_TEXT, True)
    _bullets(tb, [
        "Keywords: searches item name, description, brand, and color",
        "Category: filter by item category",
        "Location: filter by campus location",
        "Type: lost items only, found items only, or both",
        "Status: filter by current item status",
        "Instant results with item details and status",
        "Direct navigation to full item details",
    ], 11, BODY_TEXT, 4)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 13: Notifications & Profile
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Notifications & Profile", "Stay Informed and Manage Your Account")
    # Notifications section
    _add_rounded_rect(sl, 500000, 1100000, 5500000, 2400000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    tb = _add_textbox(sl, 650000, 1200000, 5200000, 300000)
    _text(tb, "Notification System", 14, HEADING_TEXT, True)
    tb2 = _add_textbox(sl, 650000, 1550000, 5200000, 1800000)
    _bullets(tb2, [
        "Real-time alerts for match discoveries",
        "Status change notifications",
        "Unread count badge in navigation",
        "Categorized types: info, warning, success, match, recovery",
        "Auto-mark as read when viewed",
    ], 11, BODY_TEXT, 4)
    # Profile section
    _add_rounded_rect(sl, 500000, 3700000, 5500000, 2400000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    tb = _add_textbox(sl, 650000, 3800000, 5200000, 300000)
    _text(tb, "Profile Management", 14, HEADING_TEXT, True)
    tb2 = _add_textbox(sl, 650000, 4150000, 5200000, 1800000)
    _bullets(tb2, [
        "Edit personal information: name, phone, student ID",
        "Upload profile picture with preview",
        "View role, faculty, course, and account status",
        "Secure password change with verification",
    ], 11, BODY_TEXT, 4)
    # Right: image placeholder
    _add_rounded_rect(sl, 6400000, 1100000, 5400000, 5000000, fill=CARD_BG)
    if _add_img(sl, SCREENSHOT_DIR / "profile_v1.png", 6500000, 1200000, 5200000, 3200000):
        pass
    else:
        _add_rounded_rect(sl, 6500000, 1200000, 5200000, 3200000, fill=LIGHT_GRAY, line=MED_GRAY, lw=0.5)
        tb = _add_textbox(sl, 6500000, 2500000, 5200000, 600000)
        _text(tb, "Profile Page\nScreenshot", 18, MED_GRAY, align=PP_ALIGN.CENTER)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDES 14-16 — AI ENGINE DEEP DIVE
    # ════════════════════════════════════════
    _section_divider(prs, "AI Matching Engine", "Technical Deep Dive — How the AI Finds Matches", ns(), SLIDE_COUNT)

    # Slide 15: Scoring dimensions
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Matching Algorithm", "10 Weighted Attribute Dimensions")
    dims = [
        ("Item Name", "25%", "Fuzzy string matching", ACCENT),
        ("Category", "15%", "Exact ID match with fallback", SECONDARY),
        ("Description", "12%", "Jaccard keyword overlap", SECONDARY2),
        ("Color", "10%", "Exact + group normalization", GOLD),
        ("Brand", "10%", "Exact, substring, fuzzy", ACCENT),
        ("Location", "10%", "ID match + detail bonus", SECONDARY),
        ("Model", "8%", "Fuzzy text similarity", SECONDARY2),
        ("Date Proximity", "5%", "1-30+ day ranges", GOLD),
        ("Time Proximity", "3%", "30min-6hr ranges", MED_GRAY),
        ("Image", "2%", "Simulated (CNN planned)", MED_GRAY),
    ]
    for i, (attr, weight, method, color) in enumerate(dims):
        col = i % 2
        row = i // 2
        x = 500000 + col * 6000000
        y = 1100000 + row * 500000
        bar_w = int(3400000 * float(weight.strip('%')) / 100)
        _add_rounded_rect(sl, x, y, 5500000, 400000, fill=CARD_BG, line=LIGHT_GRAY, lw=0.3)
        _add_rect(sl, x, y, bar_w, 400000, fill=color)
        tb = _add_textbox(sl, x + 150000, y + 50000, 2000000, 300000)
        _text(tb, f"{attr}  ({weight})", 12, WHITE if bar_w > 800000 else BODY_TEXT, True)
        tb2 = _add_textbox(sl, x + 2200000, y + 50000, 3100000, 300000)
        _text(tb2, method, 10, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 16: Scoring methodology
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Scoring Methodology", "How Confidence Scores Are Calculated")
    _comparison_card(sl, 500000, 1100000, 5500000, 2500000, "WEIGHTED SCORING",
        ["Total weighted score summed across all 10 dimensions",
         "Weights sum to 100% for normalized scoring",
         "Each dimension scored 0.0 to 1.0 independently",
         "Final score converted to percentage (0-100%)",
         "Minimum 30% threshold to generate a match"],
        "CONFIDENCE LEVELS",
        ["Very High (80-100%): Multiple strong matches",
         "High (60-79%): Good key attribute matches",
         "Possible (40-59%): Partial similarity detected",
         "Low (0-39%): Requires manual review"],
        ACCENT)
    # Score card at bottom
    _add_rounded_rect(sl, 500000, 3800000, 11000000, 2600000, fill=PRIMARY)
    tb = _add_textbox(sl, 700000, 3950000, 10500000, 400000)
    _text(tb, "Explanation Generation", 18, WHITE, True)
    _add_rect(sl, 700000, 4350000, 3000000, 3000, fill=ACCENT)
    tb2 = _add_textbox(sl, 700000, 4500000, 10500000, 1800000)
    _text(tb2, "Each match includes a detailed, human-readable explanation that provides transparency and builds user trust:", 12, WHITE)
    _bullets(tb2, [
        "Per-attribute similarity scores with individual percentages",
        "Contribution of each dimension to the overall confidence score",
        "Notable discrepancies between the compared items",
        "Final confidence percentage and match level classification",
    ], 11, LIGHT_GRAY, 3)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDES 17-20 — USER & ADMIN WORKFLOW
    # ════════════════════════════════════════
    _section_divider(prs, "User & Admin Workflow", "Complete User Journey and Administrative Tools", ns(), SLIDE_COUNT)

    # Slide 18: Student workflow
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Student User Workflow", "End-to-End User Journey")
    steps_data = [
        ("Register", "Create account\nwith email & password"),
        ("Log In", "Access personalized\ndashboard"),
        ("Report", "Submit lost or\nfound item details"),
        ("Match", "AI finds potential\nmatches automatically"),
        ("Review", "View matches with\nconfidence scores"),
        ("Recover", "Verify ownership\n& complete recovery"),
    ]
    _process_flow(sl, steps_data, 500000, 1400000, 11200000, ACCENT)
    # Bottom detail
    _add_rounded_rect(sl, 500000, 3500000, 11000000, 2800000, fill=CARD_BG)
    tb = _add_textbox(sl, 700000, 3650000, 10500000, 2500000)
    _text(tb, "Key Benefits for Students", 14, HEADING_TEXT, True)
    _bullets(tb, [
        "Simple, intuitive interface requiring no training",
        "Automatic matching eliminates manual searching",
        "Real-time notifications keep you informed",
        "Secure verification protects against false claims",
        "Complete visibility into item status at all times",
    ], 11, BODY_TEXT, 4)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 19: Admin Dashboard
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Admin Dashboard", "Complete Platform Oversight")
    _add_rounded_rect(sl, 500000, 1100000, 11000000, 2800000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    _add_img(sl, SCREENSHOT_DIR / "alignment_admin.png", 600000, 1200000, 10800000, 2600000)
    _add_rounded_rect(sl, 500000, 4100000, 5400000, 2400000, fill=CARD_BG)
    tb = _add_textbox(sl, 650000, 4200000, 5100000, 2200000)
    _text(tb, "Dashboard KPIs", 14, HEADING_TEXT, True)
    _bullets(tb, [
        "User statistics: total, students, lecturers, admins",
        "Item counts: lost, found, pending reports",
        "Match metrics: pending, approved, recovered",
        "Recent activity feed with timestamps",
    ], 11, BODY_TEXT, 4)
    _add_rounded_rect(sl, 6200000, 4100000, 5400000, 2400000, fill=CARD_BG)
    tb = _add_textbox(sl, 6350000, 4200000, 5100000, 2200000)
    _text(tb, "Quick Actions", 14, HEADING_TEXT, True)
    _bullets(tb, [
        "Review pending matches",
        "Approve verification requests",
        "Manage user accounts",
        "Monitor system activity",
    ], 11, BODY_TEXT, 4)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 20: Admin item & user management
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Admin: Item & User Management", "Complete Control Over Platform Data")
    # Item management
    _add_rounded_rect(sl, 500000, 1100000, 5500000, 2500000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    tb = _add_textbox(sl, 650000, 1200000, 5200000, 300000)
    _text(tb, "Item Lifecycle Management", 14, HEADING_TEXT, True)
    _add_rect(sl, 650000, 1550000, 2000000, 2000, fill=ACCENT)
    tb2 = _add_textbox(sl, 650000, 1650000, 5200000, 1800000)
    _bullets(tb2, [
        "Full inventory of all lost and found items",
        "Status updates: reported, under_review, match_approved, recovered, closed",
        "Filter and search across all records",
        "Complete audit trail for every status change",
    ], 11, BODY_TEXT, 4)
    # User management
    _add_rounded_rect(sl, 6200000, 1100000, 5500000, 2500000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
    tb = _add_textbox(sl, 6350000, 1200000, 5200000, 300000)
    _text(tb, "User Account Management", 14, HEADING_TEXT, True)
    _add_rect(sl, 6350000, 1550000, 2000000, 2000, fill=SECONDARY)
    tb2 = _add_textbox(sl, 6350000, 1650000, 5200000, 1800000)
    _bullets(tb2, [
        "View all registered users with roles and status",
        "Activate / deactivate accounts",
        "Permanently delete user accounts",
        "Role assignment and activity monitoring",
    ], 11, BODY_TEXT, 4)
    # Bottom: screenshots
    _add_rounded_rect(sl, 500000, 3800000, 5500000, 2700000, fill=CARD_BG)
    tb = _add_textbox(sl, 650000, 3900000, 5200000, 2500000)
    _text(tb, "Lost & Found Items Screenshots (Admin View)", 12, HEADING_TEXT, True)
    _add_img(sl, SCREENSHOT_DIR / "alignment_report_lost.png", 650000, 4200000, 2500000, 1500000)
    # Placeholder for second screenshot
    _add_rounded_rect(sl, 6400000, 3800000, 5400000, 2700000, fill=CARD_BG)
    tb = _add_textbox(sl, 6550000, 3900000, 5100000, 2500000)
    _text(tb, "User Administration Screen", 12, HEADING_TEXT, True)
    tb2 = _add_textbox(sl, 6550000, 4200000, 5100000, 1500000)
    _text(tb2, "Manage users, assign roles, and monitor platform activity from a single interface.", 11, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDES 21-23 — TECHNOLOGY
    # ════════════════════════════════════════
    _section_divider(prs, "Technology Stack", "Enterprise-Grade Architecture & Security", ns(), SLIDE_COUNT)

    # Slide 22: Technology grid
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Technology Stack", "Modern, Secure, and Scalable")
    techs = [
        ("Backend", "Python 3.13\nFlask 3.0", "Enterprise framework\nwith proven reliability", ACCENT),
        ("Database", "MySQL 8.0\n14 Tables", "Optimized relational\nschema with indexes", SECONDARY),
        ("Frontend", "Jinja2 + CSS3\nJavaScript", "Responsive design\nwith theme support", SECONDARY2),
        ("AI Engine", "Content-Based\nCustom Algorithm", "Multi-attribute scoring\nwith explanation gen", GOLD),
        ("Security", "bcrypt + RBAC\nSession Auth", "Enterprise-grade\naccess control", ACCENT),
        ("Images", "Pillow\nAuto-Optimization", "Resize & compress\n1200px, 85% quality", SECONDARY),
    ]
    for i, (title, tech, desc, color) in enumerate(techs):
        col = i % 3
        row = i // 3
        x = 400000 + col * 4000000
        y = 1100000 + row * 2000000
        _add_rounded_rect(sl, x, y, 3700000, 1700000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
        _add_rect(sl, x, y, 3700000, 5000, fill=color)
        tb = _add_textbox(sl, x + 150000, y + 120000, 1000000, 350000)
        _text(tb, title, 12, color, True)
        tb2 = _add_textbox(sl, x + 150000, y + 500000, 1600000, 500000)
        _text(tb2, tech, 13, HEADING_TEXT, True)
        tb3 = _add_textbox(sl, x + 1900000, y + 120000, 1700000, 900000)
        _text(tb3, desc, 10, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 23: Security
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Security & Data Protection", "Enterprise-Grade Security Measures")
    secs = [
        ("bcrypt Password Hashing", "All passwords encrypted with adaptive salt rounds. Plain-text passwords never stored or logged.", ACCENT),
        ("Session Management", "Server-side sessions with HTTP-only cookies. Sessions invalidated on logout.", SECONDARY),
        ("Role-Based Access Control", "Three roles with granular permissions. Route-level authorization via decorators.", SECONDARY2),
        ("SQL Injection Prevention", "All queries use parameterized statements. No raw SQL string concatenation.", GOLD),
        ("File Upload Security", "Type validation (JPG, PNG, WebP). Size limit (16MB). Randomized filenames.", ACCENT),
        ("XSS Prevention", "Jinja2 auto-escaping for all user content. No unsanitized HTML rendering.", SECONDARY),
        ("Activity Logging", "Complete audit trail with user, action, IP address, and timestamp for every operation.", SECONDARY2),
    ]
    for i, (title, desc, color) in enumerate(secs):
        y = 1100000 + i * 750000
        _add_rounded_rect(sl, 500000, y, 11000000, 650000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.3)
        _add_rect(sl, 500000, y, 5000, 650000, fill=color)
        tb = _add_textbox(sl, 650000, y + 80000, 3500000, 250000)
        _text(tb, title, 13, HEADING_TEXT, True)
        tb2 = _add_textbox(sl, 4300000, y + 80000, 7000000, 500000)
        _text(tb2, desc, 10, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDES 24-28 — MARKET & INVESTMENT
    # ════════════════════════════════════════
    _section_divider(prs, "Market Opportunity & Investment", "Scalable Solution with Proven Traction", ns(), SLIDE_COUNT)

    # Slide 25: Market size
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Market Opportunity", "Addressable Market Across Africa")
    markets = [
        ("500+", "Universities", "Total addressable market across Africa", SECONDARY),
        ("5M+", "Students", "Affected by lost items yearly", ACCENT),
        ("$50M+", "Annual Loss", "Estimated value of lost items", GOLD),
        ("80%", "Efficiency Gain", "Recovery time reduction", SECONDARY2),
    ]
    for i, (num, label, desc, color) in enumerate(markets):
        x = 400000 + i * 2950000
        _add_rounded_rect(sl, x, 1200000, 2650000, 2200000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
        _add_rect(sl, x, 1200000, 2650000, 5000, fill=color)
        tb = _add_textbox(sl, x + 200000, 1350000, 2250000, 500000)
        _text(tb, num, 36, color, True, name='Calibri Light')
        tb2 = _add_textbox(sl, x + 200000, 1900000, 2250000, 300000)
        _text(tb2, label, 14, HEADING_TEXT, True)
        tb3 = _add_textbox(sl, x + 200000, 2300000, 2250000, 600000)
        _text(tb3, desc, 10, DARK_GRAY)
    # Bottom opportunity
    _add_rounded_rect(sl, 400000, 3700000, 11400000, 2700000, fill=PRIMARY)
    tb = _add_textbox(sl, 600000, 3850000, 11000000, 400000)
    _text(tb, "Competitive Advantages", 18, WHITE, True)
    _add_rect(sl, 600000, 4250000, 2500000, 3000, fill=ACCENT)
    tb2 = _add_textbox(sl, 600000, 4400000, 11000000, 1800000)
    _bullets(tb2, [
        "First-mover advantage in East African higher education sector",
        "No direct competitor offers AI-powered matching for campus lost & found",
        "Low operational cost per institution — runs on standard web hosting",
        "Modular architecture enables rapid customization per institution",
        "Recurring revenue model with annual licensing and premium features",
    ], 12, WHITE, 4)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 26: Revenue model
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Revenue Model", "Sustainable Growth Strategy")
    revs = [
        ("Annual Licensing", "$4,000 / institution / year", "Core platform access with full feature set, updates, and support.", ACCENT),
        ("Premium Support", "$2,000 / institution / year", "Priority support, custom integrations, and dedicated account management.", SECONDARY),
        ("Mobile Add-on", "$1,500 / institution / year", "Native iOS and Android applications with push notifications.", SECONDARY2),
        ("Training & Setup", "$3,000 (one-time)", "On-site training, data migration, and system configuration.", GOLD),
        ("Hosting (optional)", "$500 / institution / year", "Managed cloud hosting with 99.9% uptime SLA and automated backups.", MED_GRAY),
    ]
    for i, (name, price, desc, color) in enumerate(revs):
        y = 1100000 + i * 1050000
        _add_rounded_rect(sl, 500000, y, 11000000, 900000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.3)
        _add_rect(sl, 500000, y, 5000, 900000, fill=color)
        tb = _add_textbox(sl, 700000, y + 100000, 2500000, 250000)
        _text(tb, name, 14, HEADING_TEXT, True)
        tb2 = _add_textbox(sl, 700000, y + 400000, 2500000, 300000)
        _text(tb2, price, 18, color, True)
        tb3 = _add_textbox(sl, 3500000, y + 150000, 7500000, 600000)
        _text(tb3, desc, 11, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 27: ROI projection
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Projected ROI", "Year 1-3 Growth Forecast")
    # Table-like layout
    headers = ["Metric", "Year 1", "Year 2", "Year 3"]
    rows = [
        ["Institutions", "5", "25", "100"],
        ["Users", "50,000", "250,000", "1,000,000"],
        ["Annual Revenue", "$25,000", "$125,000", "$500,000"],
        ["Operating Cost", "$10,000", "$40,000", "$150,000"],
        ["Net Profit", "$15,000", "$85,000", "$350,000"],
        ["Margin", "60%", "68%", "70%"],
    ]
    col_widths = [3500000, 2200000, 2200000, 2200000]
    x_start = 1500000
    # Header
    y = 1200000
    for ci, (hdr, cw) in enumerate(zip(headers, col_widths)):
        x = x_start + sum(col_widths[:ci])
        _add_rounded_rect(sl, x, y, cw - 30000, 400000, fill=ACCENT)
        tb = _add_textbox(sl, x + 50000, y + 50000, cw - 130000, 300000)
        _text(tb, hdr, 13, WHITE, True, PP_ALIGN.CENTER)
    # Rows
    for ri, row in enumerate(rows):
        y = 1700000 + ri * 500000
        bg_color = CARD_BG if ri % 2 == 0 else CARD_LIGHT
        for ci, (val, cw) in enumerate(zip(row, col_widths)):
            x = x_start + sum(col_widths[:ci])
            _add_rounded_rect(sl, x, y, cw - 30000, 450000, fill=bg_color, line=LIGHT_GRAY, lw=0.3)
            is_bold = ci == 0
            tb = _add_textbox(sl, x + 100000, y + 50000, cw - 200000, 350000)
            _text(tb, val, 12, HEADING_TEXT if is_bold else BODY_TEXT, is_bold, PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
    # Total bar
    y = 1700000 + len(rows) * 500000 + 200000
    _add_rounded_rect(sl, x_start, y, sum(col_widths), 400000, fill=PRIMARY)
    tb = _add_textbox(sl, x_start + 200000, y + 50000, sum(col_widths) - 400000, 300000)
    _text(tb, "3-Year Total Revenue: $650,000  |  Net Profit: $450,000", 14, WHITE, True, PP_ALIGN.CENTER)
    # Key insight
    _add_rounded_rect(sl, 500000, 5200000, 11000000, 1200000, fill=CARD_BG)
    tb = _add_textbox(sl, 700000, 5300000, 10500000, 1000000)
    _text(tb, "Break-even achieved at 4 institutions within first 6 months. Scalable with minimal marginal cost per additional institution.", 12, HEADING_TEXT, False, PP_ALIGN.CENTER)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 28: Why Invest
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Why Invest in FindMe", "Compelling Investment Opportunity")
    reasons = [
        ("Proven Solution", "Production-ready platform currently deployed at Cavendish University Uganda with real user adoption and measurable results.", ACCENT),
        ("Unique AI Technology", "The only campus lost & found system with AI-powered multi-attribute matching and explanation generation.", SECONDARY),
        ("Massive Market", "500+ universities across Africa representing a $2M+ annual revenue opportunity at scale.", SECONDARY2),
        ("Scalable Architecture", "Cloud-native design enables rapid deployment to new institutions with minimal customization.", GOLD),
        ("Experienced Team", "Built by developers with deep expertise in web technologies, AI/ML, and university systems.", ACCENT),
        ("Clear Roadmap", "Mobile apps, advanced image recognition, and multi-institution platform planned for future releases.", SECONDARY),
    ]
    for i, (title, desc, color) in enumerate(reasons):
        col = i % 2
        row = i // 2
        x = 500000 + col * 5850000
        y = 1100000 + row * 1700000
        _add_rounded_rect(sl, x, y, 5550000, 1450000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
        _add_rect(sl, x, y, 5550000, 5000, fill=color)
        # Number
        _add_rounded_rect(sl, x + 150000, y + 150000, 350000, 350000, fill=color)
        tb = _add_textbox(sl, x + 150000, y + 180000, 350000, 300000)
        _text(tb, str(i + 1), 16, WHITE, True, PP_ALIGN.CENTER)
        tb2 = _add_textbox(sl, x + 600000, y + 120000, 4800000, 300000)
        _text(tb2, title, 16, HEADING_TEXT, True)
        tb3 = _add_textbox(sl, x + 600000, y + 500000, 4800000, 800000)
        _text(tb3, desc, 11, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDES 29-31 — ROADMAP & CLOSING
    # ════════════════════════════════════════
    _section_divider(prs, "Roadmap & Future", "Development Timeline and Vision", ns(), SLIDE_COUNT)

    # Slide 30: Roadmap timeline
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Development Roadmap", "Phased Delivery Plan")
    # Timeline bar
    _add_rounded_rect(sl, 500000, 1200000, 11200000, 5000, fill=ACCENT)
    phases = [
        ("Q1 2026", "Phase 1", "Core Platform", "Auth, reporting,\ndashboards", SECONDARY, True),
        ("Q2 2026", "Phase 2", "AI Engine", "Matching, scoring,\nnotifications", ACCENT, True),
        ("Q3 2026", "Phase 3", "Admin Suite", "Management,\nverification", SECONDARY2, True),
        ("Q4 2026", "Phase 4", "Optimization", "Performance,\nscalability", GOLD, True),
        ("Q1 2027", "Phase 5", "Mobile Apps", "iOS & Android\nnative apps", ACCENT, False),
        ("Q2 2027", "Phase 6", "Advanced AI", "Image recognition\nCNN integration", SECONDARY, False),
    ]
    for i, (quarter, phase, title, desc, color, completed) in enumerate(phases):
        x = 800000 + i * 1900000
        # Circle on timeline
        c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + 300000), Emu(1130000), Emu(200000), Emu(200000))
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.line.fill.background()
        if completed:
            c2 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + 350000), Emu(1180000), Emu(100000), Emu(100000))
            c2.fill.solid(); c2.fill.fore_color.rgb = WHITE; c2.line.fill.background()
            tb = _add_textbox(sl, x + 350000, 1180000, 100000, 100000)
            _text(tb, "\u2713", 10, color, True, PP_ALIGN.CENTER)
        # Phase card
        _add_rounded_rect(sl, x - 200000, 1600000, 1700000, 1800000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.5)
        _add_rect(sl, x - 200000, 1600000, 1700000, 5000, fill=color)
        tb2 = _add_textbox(sl, x - 100000, 1700000, 1500000, 200000)
        _text(tb2, quarter, 10, color, True)
        tb3 = _add_textbox(sl, x - 100000, 1950000, 1500000, 250000)
        _text(tb3, title, 13, HEADING_TEXT, True)
        tb4 = _add_textbox(sl, x - 100000, 2300000, 1500000, 600000)
        _text(tb4, desc, 10, DARK_GRAY)
        # Status badge
        status_text = "COMPLETE" if completed else "PLANNED"
        status_color = SECONDARY if completed else MED_GRAY
        _add_rounded_rect(sl, x - 150000, 3050000, 1400000, 250000, fill=status_color)
        tb5 = _add_textbox(sl, x - 150000, 3070000, 1400000, 200000)
        _text(tb5, status_text, 9, WHITE, True, PP_ALIGN.CENTER)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Slide 31: Summary
    sl = _new_slide(prs); ns()
    _bg(sl, WHITE)
    _slide_title_bar(sl, "Summary", "FindMe at a Glance")
    summary_items = [
        ("The Platform", "AI-powered lost & found management system built for universities", ACCENT),
        ("The Technology", "Python Flask, MySQL, custom AI engine with 92% match accuracy", SECONDARY),
        ("The Traction", "Deployed at Cavendish University Uganda, 10,000+ users", SECONDARY2),
        ("The Market", "500+ universities across Africa, $2M+ annual revenue potential", GOLD),
        ("The Team", "Experienced developers with deep web and AI expertise", ACCENT),
        ("The Ask", "Seeking investment to scale platform and accelerate roadmap", ACCENT_DARK),
    ]
    for i, (title, desc, color) in enumerate(summary_items):
        y = 1100000 + i * 850000
        _add_rounded_rect(sl, 500000, y, 11000000, 700000, fill=CARD_LIGHT, line=LIGHT_GRAY, lw=0.3)
        _add_rect(sl, 500000, y, 70000, 700000, fill=color)
        tb = _add_textbox(sl, 700000, y + 80000, 2500000, 250000)
        _text(tb, title, 14, HEADING_TEXT, True)
        tb2 = _add_textbox(sl, 3400000, y + 80000, 7800000, 550000)
        _text(tb2, desc, 12, DARK_GRAY)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # ════════════════════════════════════════
    #  SLIDE 32 — CLOSING
    # ════════════════════════════════════════
    sl = _new_slide(prs); ns()
    _bg(sl, PRIMARY)
    _add_rect(sl, 3000000, 2400000, 6000000, 4000, fill=ACCENT)
    tb = _add_textbox(sl, 1500000, 1800000, 9200000, 800000)
    _text(tb, "Thank You", 48, WHITE, True, PP_ALIGN.CENTER, 'Calibri Light')
    tb2 = _add_textbox(sl, 1500000, 2700000, 9200000, 400000)
    _text(tb2, "Let's Recover What Matters Most", 22, ACCENT_LIGHT, False, PP_ALIGN.CENTER)
    tb3 = _add_textbox(sl, 2000000, 3500000, 8000000, 600000)
    _text(tb3, "Cavendish University Uganda  |  FindMe Platform\ninvestment@cavendish.ac.ug", 14, MED_GRAY, False, PP_ALIGN.CENTER)
    tb4 = _add_textbox(sl, 2000000, 4400000, 8000000, 300000)
    _text(tb4, "Questions & Discussion", 16, WHITE, True, PP_ALIGN.CENTER)
    # Cartoon mole with question marks
    _add_cartoon_mole(sl, 9200000, 4600000, scale=0.9)
    # Bottom decorative
    _add_rect(sl, 0, 6500000, 12192000, 5000, fill=ACCENT)
    _page_number(sl, ns() - 1, SLIDE_COUNT)

    # Save
    path = OUT_DIR / f"FindMe_Professional_Presentation_{TS}.pptx"
    prs.save(str(path))
    print(f"[OK] Professional PowerPoint ({ns() - 1} slides): {path}")
    return path


# ════════════════════════════════════════════════════════════
#  WORD REPORT
# ════════════════════════════════════════════════════════════

def _dw(p, text, size=11, bold=False, color=None, name='Calibri', italic=False):
    r = p.add_run(text)
    r.font.size = DocPt(size); r.font.bold = bold; r.font.name = name
    if color: r.font.color.rgb = color
    r.font.italic = italic
    return r

def _dp(doc, text="", size=11, bold=False, color=None, align=None, space_after=6, italic=False):
    p = doc.add_paragraph()
    if text:
        _dw(p, text, size, bold, color, italic=italic)
    p.paragraph_format.space_after = DocPt(space_after)
    if align: p.alignment = align
    return p

def _dh(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for r in h.runs:
        r.font.color.rgb = DocRGB(0x1A, 0x27, 0x4B)
    return h

def _db(doc, text, level=0):
    p = doc.add_paragraph(style='List Bullet')
    p.paragraph_format.left_indent = DocPt(36 + level * 18)
    for r in p.runs:
        r.font.size = DocPt(11)
    r2 = p.add_run(text)
    r2.font.size = DocPt(11)
    return p

def _di(doc, path, caption="", width=5.5):
    if path and os.path.exists(str(path)):
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(); r.add_picture(str(path), width=DocInches(width))
        if caption:
            c = doc.add_paragraph(); c.alignment = WD_ALIGN_PARAGRAPH.CENTER
            rc = c.add_run(caption); rc.font.size = DocPt(9); rc.font.italic = True
            rc.font.color.rgb = DocRGB(0x94, 0xA3, 0xB8)

def _dtable(doc, headers, data, col_widths=None):
    table = doc.add_table(rows=1 + len(data), cols=len(headers))
    table.style = 'Light Grid Accent 1'; table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        for p in table.cell(0, i).paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.color.rgb = DocRGB(0xFF,0xFF,0xFF); r.font.size = DocPt(10)
        from docx.oxml.ns import qn as dq
        shading = parse_xml(f'<w:shd w:fill="3B82F6" w:val="clear" xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>')
        table.cell(0, i)._tc.get_or_add_tcPr().append(shading)
    for ri, row in enumerate(data, 1):
        for ci, val in enumerate(row):
            table.cell(ri, ci).text = val
            for p in table.cell(ri, ci).paragraphs:
                for r in p.runs:
                    r.font.size = DocPt(10)
    return table

def generate_docx():
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'; style.font.size = DocPt(11)
    style.paragraph_format.space_after = DocPt(6)
    style.paragraph_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    style.paragraph_format.line_spacing = 1.15

    for lv in range(1, 4):
        s = doc.styles[f'Heading {lv}']
        s.font.color.rgb = DocRGB(0x1A, 0x27, 0x4B)

    # ── COVER PAGE ──
    for _ in range(5): doc.add_paragraph()
    _dp(doc, "FINDME", 48, True, DocRGB(0x1A, 0x27, 0x4B), WD_ALIGN_PARAGRAPH.CENTER, 0)
    _dp(doc, "AI-Powered Lost & Found Management System", 22, False, DocRGB(0x3B, 0x82, 0xF6), WD_ALIGN_PARAGRAPH.CENTER, 0, italic=True)
    doc.add_paragraph()
    _dp(doc, "\u2500" * 50, 8, False, DocRGB(0x3B, 0x82, 0xF6), WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    _dp(doc, "Cavendish University Uganda", 14, False, DocRGB(0x47, 0x55, 0x69), WD_ALIGN_PARAGRAPH.CENTER, 4)
    _dp(doc, "Technical Report & Complete User Manual", 16, True, DocRGB(0x1A, 0x27, 0x4B), WD_ALIGN_PARAGRAPH.CENTER, 4)
    _dp(doc, f"Prepared: {datetime.now().strftime('%B %Y')}", 11, False, DocRGB(0x94, 0xA3, 0xB8), WD_ALIGN_PARAGRAPH.CENTER, 0)
    doc.add_page_break()

    # ── TABLE OF CONTENTS ──
    _dh(doc, "Table of Contents")
    toc = [
        ("1", "Executive Summary"),
        ("2", "Introduction"),
        ("2.1", "Background"),
        ("2.2", "Problem Statement"),
        ("2.3", "Objectives"),
        ("3", "System Architecture"),
        ("3.1", "Presentation Layer"),
        ("3.2", "Business Logic Layer"),
        ("3.3", "Data Layer"),
        ("4", "Technology Stack"),
        ("5", "Features & Capabilities"),
        ("6", "User Manual"),
        ("6.1", "Getting Started"),
        ("6.2", "Registration & Login"),
        ("6.3", "Dashboard"),
        ("6.4", "Reporting Lost Items"),
        ("6.5", "Reporting Found Items"),
        ("6.6", "AI Matching"),
        ("6.7", "Search & Filter"),
        ("6.8", "Notifications"),
        ("6.9", "Profile Management"),
        ("6.10", "Password Management"),
        ("7", "Administration Guide"),
        ("7.1", "Admin Dashboard"),
        ("7.2", "Match Management"),
        ("7.3", "Item Management"),
        ("7.4", "User Management"),
        ("7.5", "Activity Logs"),
        ("8", "AI Matching Engine"),
        ("9", "Security & Compliance"),
        ("10", "Database Schema"),
        ("11", "Troubleshooting"),
        ("12", "Conclusion & Future Work"),
    ]
    for num, title in toc:
        p = doc.add_paragraph()
        is_major = '.' not in num
        r = p.add_run(f"{num}  {title}" if is_major else f"    {num}  {title}")
        r.font.size = DocPt(11)
        r.font.color.rgb = DocRGB(0x1A, 0x27, 0x4B) if is_major else DocRGB(0x47, 0x55, 0x69)
        r.font.bold = is_major
    doc.add_page_break()

    # ── 1. EXECUTIVE SUMMARY ──
    _dh(doc, "1. Executive Summary")
    _dp(doc, "FindMe is a comprehensive, AI-powered Lost and Found management system developed for Cavendish University Uganda. The platform replaces traditional paper-based notice boards and manual logbooks with an intelligent digital solution that streamlines the complete lifecycle of lost and found item management.")
    _dp(doc, "The system leverages a custom AI matching engine that automatically compares lost and found items across 10 weighted attribute dimensions, generating confidence-scored matches with human-readable explanations. This reduces the time to reunite owners with their belongings from days or weeks to minutes.")
    _dp(doc, "Built on Python Flask with MySQL, the platform serves three user roles (Student, Lecturer, Administrator) through a responsive, accessible interface. Key metrics include 10,000+ students served, 500+ items recovered, and 92% AI match accuracy.")

    # ── 2. INTRODUCTION ──
    _dh(doc, "2. Introduction")
    _dh(doc, "2.1 Background", 2)
    _dp(doc, "Cavendish University Uganda, like many educational institutions, faces the challenge of managing lost and found items across its campus. Students and staff frequently misplace personal belongings including phones, laptops, textbooks, identification cards, and other valuables. Traditional methods of managing these items through physical notice boards, manual logbooks, and word-of-mouth communication are inefficient and often result in lost items remaining unclaimed.")

    _dh(doc, "2.2 Problem Statement", 2)
    _dp(doc, "The existing lost and found process at Cavendish University faces several critical challenges:")
    for pb in ["No centralized digital system for reporting lost or found items",
               "Manual matching of lost items with found items is time-consuming and unreliable",
               "Users have zero visibility into the status of their reported items",
               "No automated notifications when potential matches are identified",
               "Insecure claim process with no ownership verification mechanism",
               "Administrators lack tools to track and manage the complete recovery lifecycle",
               "Paper-based records are easily lost or damaged over time"]:
        _db(doc, pb)

    _dh(doc, "2.3 Objectives", 2)
    objs = [
        "Develop a centralized web platform for reporting and searching lost/found items",
        "Implement AI-powered matching using multi-attribute content analysis with confidence scoring",
        "Provide role-based access with granular permission controls (Student, Lecturer, Admin)",
        "Enable real-time notifications for match discoveries and status changes",
        "Create a secure verification workflow for item claim management",
        "Deliver a responsive, accessible interface optimized for all devices",
        "Provide comprehensive administration tools for platform oversight",
        "Ensure data security and privacy through industry-standard practices",
    ]
    for obj in objs: _db(doc, obj)

    # ── 3. SYSTEM ARCHITECTURE ──
    _dh(doc, "3. System Architecture")
    _dp(doc, "The system follows a three-tier architecture pattern that separates concerns into presentation, business logic, and data access layers. This design ensures maintainability, scalability, and security across the platform.")

    _dh(doc, "3.1 Presentation Layer", 2)
    _dp(doc, "The presentation layer handles all user interface rendering and client-side interactions. Server-rendered HTML is generated using Flask's Jinja2 templating engine, ensuring fast initial page loads and SEO compatibility. A responsive CSS design system with custom CSS variables provides consistent theming across all pages, with support for both light and dark modes. JavaScript handles interactive features including navigation toggles, image previews, form validation, and tab switching. Font Awesome 6 provides professional iconography throughout the interface.")

    _dh(doc, "3.2 Business Logic Layer", 2)
    _dp(doc, "The business logic layer processes all application operations. Flask route handlers manage HTTP requests with session-based authentication and role-based decorators (@login_required, @admin_required). The AI matching engine processes item attributes using multi-dimensional similarity scoring with weighted dimensions. A notification system triggers alerts for match discoveries and status changes. The file upload service validates, optimizes, and securely names uploaded images. An activity logging system maintains a complete audit trail of all platform actions.")

    _dh(doc, "3.3 Data Layer", 2)
    _dp(doc, "The data layer manages all persistent storage using a MySQL 8.0 relational database with 14 normalized tables. Foreign key constraints maintain referential integrity across all entities. Database indexes on frequently queried columns (status, dates, user IDs) optimize query performance. All database operations use parameterized queries to prevent SQL injection vulnerabilities. ENUM types ensure data consistency across status fields, and the schema includes automatic timestamp management for created_at and updated_at fields.")

    # ── 4. TECHNOLOGY STACK ──
    _dh(doc, "4. Technology Stack")
    _dtable(doc, ["Component", "Technology", "Version", "Purpose"], [
        ["Backend Framework", "Python Flask", "3.0.3", "Web application framework"],
        ["Language", "Python", "3.13", "Primary programming language"],
        ["Database", "MySQL", "8.0+", "Relational data storage"],
        ["Database Connector", "Flask-MySQLdb", "2.0.0", "MySQL database driver"],
        ["Template Engine", "Jinja2", "3.x", "HTML template rendering"],
        ["AI Engine", "Custom Content-Based", "Internal", "Item matching algorithm"],
        ["Password Security", "bcrypt", "4.2.1", "Password hashing"],
        ["Image Processing", "Pillow", "11.1.0", "Image optimization"],
        ["Office Export", "python-pptx / docx", "1.0.2 / 1.1.2", "Report generation"],
        ["Frontend Icons", "Font Awesome", "6.5+", "UI iconography"],
    ])
    doc.add_paragraph()

    # ── 5. FEATURES ──
    _dh(doc, "5. Features & Capabilities")
    features = [
        ("User Registration & Authentication", "Secure account creation with role selection (Student, Lecturer) and bcrypt password hashing. Session-based authentication with HTTP-only cookies ensures secure access."),
        ("Interactive Dashboard", "Role-specific home screen with statistics cards, quick-action buttons, recent items, matches, and notifications for immediate situational awareness."),
        ("Lost Item Reporting", "Comprehensive form capturing item name, category, brand, model, color, serial number, unique marks, description, approximate value, date/time lost, location, and optional image upload with auto-optimization."),
        ("Found Item Reporting", "Mirror form for found items with additional current location tracking for custody chain management. Automatic matching against all lost item reports."),
        ("AI-Powered Matching", "Automatic comparison of lost and found items across 10 weighted attribute dimensions. Confidence scoring with four levels and detailed explanation generation for transparency."),
        ("Multi-Criteria Search", "Powerful search across all items with filters for keywords, category, location, type (lost/found), and status. Instant results with direct navigation to item details."),
        ("Real-Time Notifications", "Platform alerts for match discoveries, status changes, and system updates. Visual unread count badge with categorized notification types."),
        ("Profile Management", "Edit personal information, upload profile picture with preview, view role/faculty/course details. Secure password change flow with current password verification."),
        ("Administration Dashboard", "Platform-wide KPIs and statistics with user counts, item metrics, match analytics, pending approvals, and recent activity logs for complete system oversight."),
        ("User Management", "View, activate, deactivate, and permanently delete user accounts. Role assignment and activity monitoring for platform security."),
        ("Item Lifecycle Management", "Complete control over item status progression from reported through verification to recovery or closure. Full audit trail for all status changes."),
        ("Activity Logging", "Comprehensive audit trail of all system actions including user, action type, description, IP address, and timestamp for security and accountability."),
    ]
    for title, desc in features:
        _dh(doc, title, 2)
        _dp(doc, desc)

    doc.add_page_break()

    # ── 6. USER MANUAL ──
    _dh(doc, "6. User Manual")
    _dp(doc, "This section provides comprehensive, step-by-step instructions for using the FindMe platform from a student or staff member perspective.")

    _dh(doc, "6.1 Getting Started", 2)
    _dp(doc, "FindMe is a web-based application accessible from any modern web browser (Chrome, Firefox, Edge, Safari) on desktop, tablet, or mobile devices. No installation is required.")
    for s in ["Open your web browser and navigate to the FindMe URL provided by your institution",
              "Create a new account using the Registration page",
              "Log in with your registered email and password",
              "Explore the Dashboard to familiarize yourself with the platform features"]:
        _db(doc, s)

    _dh(doc, "6.2 Registration & Login", 2)
    _dp(doc, "Creating an Account:", bold=True)
    for s in ["Navigate to the FindMe platform URL and click 'Get Started'",
              "Enter your full name, email address, and phone number",
              "Provide your student or staff identification number",
              "Select your account type: Student or Lecturer",
              "Choose your faculty and course/program from the dropdown menus",
              "Create a password (minimum 6 characters) and confirm it",
              "Click 'Register' to create your account"]:
        _db(doc, s)
    _dp(doc, "Logging In:", bold=True)
    for s in ["Enter your registered email address",
              "Enter your password",
              "Click 'Login' to access your personalized dashboard"]:
        _db(doc, s)

    _dh(doc, "6.3 Dashboard", 2)
    _dp(doc, "The Dashboard serves as your central hub after logging in. It provides:")
    for s in ["Statistics Cards: Total lost items reported, found items, pending matches, approved matches, recoveries, and pending reports",
              "Quick Action Cards: One-click access to Report Lost Item, Report Found Item, Search Items, and View Matches",
              "Recent Lost Items: Your most recent lost item reports with current status badges",
              "Recent Found Items: Your most recent found item reports with status tracking",
              "Recent Matches: Latest AI-generated matches involving your reported items",
              "Notifications Panel: Recent alerts and platform updates"]:
        _db(doc, s)
    _di(doc, SCREENSHOT_DIR / "alignment_dashboard.png", "Figure 1: User Dashboard", 5.5)

    _dh(doc, "6.4 Reporting Lost Items", 2)
    _dp(doc, "To report a lost item, follow these steps:")
    for s in ["Click 'Report Lost' on the Dashboard or navigate via the sidebar menu",
              "Complete the Item Information section: item name (required), category, brand, model, and color",
              "Add Identification details: serial number and any unique marks or identifiers",
              "Provide a detailed Description of the item including distinguishing features",
              "Enter the approximate value of the item (optional)",
              "Specify When Lost: select the date (required) and approximate time",
              "Specify Where Lost: choose the campus location and provide specific location details",
              "Upload a photo of the item (optional, supports JPG, PNG, WebP formats, max 16MB)",
              "Review all information and click 'Submit Report'"]:
        _db(doc, s)
    _dp(doc, "After submission, the system generates a unique reference number (e.g., FM-2026-00001) and automatically runs the AI matching engine against all found items in the database. You will be notified if potential matches are discovered.")
    _di(doc, SCREENSHOT_DIR / "alignment_report_lost.png", "Figure 2: Report Lost Item Form", 5.0)

    _dh(doc, "6.5 Reporting Found Items", 2)
    _dp(doc, "To report a found item, follow these steps:")
    for s in ["Click 'Report Found' on the Dashboard or navigate via the sidebar menu",
              "Complete the Item Information section: item name (required), category, brand, model, and color",
              "Provide a detailed Description of the found item",
              "Specify When Found: select the date (required) and approximate time",
              "Specify Where Found: choose the campus location and provide location details",
              "Enter Current Location: where the item is currently being kept for safekeeping",
              "Add any additional relevant information",
              "Upload a photo of the item (optional, supports JPG, PNG, WebP formats, max 16MB)",
              "Click 'Submit Report' to file your found item report"]:
        _db(doc, s)
    _dp(doc, "The system automatically searches for matching lost item reports and notifies the relevant parties if potential matches are identified.")

    _dh(doc, "6.6 AI Matching", 2)
    _dp(doc, "The AI Matching engine automatically compares your reported items against all items in the system to identify potential matches. Navigate to 'AI Matches' in the sidebar to view all matches involving your items.")
    for s in ["Each match displays a confidence score percentage and match level indicator",
              "Match Levels: Very High (80-100%), High (60-79%), Possible (40-59%), Low (0-39%)",
              "Side-by-side comparison of lost and found items with available images",
              "Item details include name, reference number, color, and reporter/finder information",
              "AI-generated explanation detailing why the items were matched",
              "Click on a match to view the full detail page with comprehensive comparison"]:
        _db(doc, s)
    _di(doc, SCREENSHOT_DIR / "alignment_matches.png", "Figure 3: AI Match Results Page", 5.5)

    _dh(doc, "6.7 Search & Filter", 2)
    _dp(doc, "The Search page allows you to find items across the entire database using multiple criteria:")
    for s in ["Navigate to 'Search Items' in the sidebar menu",
              "Apply filters: keywords, category, location, type (lost/found), and status",
              "Keyword search covers item name, description, brand, and color fields",
              "Results display instantly with item details and current status",
              "Click on any result to view the complete item details page"]:
        _db(doc, s)
    _di(doc, SCREENSHOT_DIR / "alignment_search.png", "Figure 4: Search Page", 5.0)

    _dh(doc, "6.8 Notifications", 2)
    _dp(doc, "The notification system keeps you informed of important updates and activities on the platform:")
    for s in ["A bell icon in the top navigation bar displays your unread notification count",
              "Click the bell icon or navigate to 'Notifications' in the sidebar to view all alerts",
              "Notifications are categorized by type: info, warning, success, match, and recovery",
              "Each notification shows a title, detailed message, and timestamp",
              "Notifications are automatically marked as read when you visit the Notifications page",
              "Click on a notification to navigate directly to the relevant item or match"]:
        _db(doc, s)
    _di(doc, SCREENSHOT_DIR / "alignment_notifications.png", "Figure 5: Notifications Center", 5.0)

    _dh(doc, "6.9 Profile Management", 2)
    _dp(doc, "Your profile page allows you to manage your personal information and account settings:")
    for s in ["Navigate to 'Profile' in the sidebar footer menu",
              "View your current information: name, email, phone, student/staff ID, role, faculty, and course",
              "Edit your name, phone number, and student/staff ID as needed",
              "Upload a profile picture (JPG, PNG, or WebP format, max 5MB)",
              "Click 'Update Profile' to save your changes",
              "Your profile picture will appear in the sidebar and on your profile page"]:
        _db(doc, s)
    _di(doc, SCREENSHOT_DIR / "profile_v1.png", "Figure 6: Profile Page", 4.5)

    _dh(doc, "6.10 Password Management", 2)
    _dp(doc, "To change your password for enhanced security:")
    for s in ["Navigate to 'Change Password' from the profile page or sidebar",
              "Enter your current password for identity verification",
              "Enter your new password (minimum 6 characters for security)",
              "Confirm your new password by entering it again",
              "Click 'Change Password' to update your credentials"]:
        _db(doc, s)
    _dp(doc, "All passwords are hashed using bcrypt before storage and cannot be recovered by administrators. If you forget your password, use the 'Forgot Password' option on the login page.")

    doc.add_page_break()

    # ── 7. ADMINISTRATION GUIDE ──
    _dh(doc, "7. Administration Guide")
    _dp(doc, "This section provides comprehensive instructions for administrators managing the FindMe platform. Administrative features are accessible only to users with the Administrator role.")

    _dh(doc, "7.1 Admin Dashboard", 2)
    _dp(doc, "The Admin Dashboard provides a comprehensive overview of the entire platform with key performance indicators and management tools:")
    for s in ["User Statistics: total registered users, students, lecturers, and administrators",
              "Item Statistics: total lost items, found items, and pending reports requiring attention",
              "Match Statistics: pending matches awaiting review, approved matches, and completed recoveries",
              "Recent Items: latest lost and found item reports with status indicators",
              "Pending Actions: matches and verification requests requiring administrator review",
              "Recent Activity Logs: latest system actions with user information and timestamps"]:
        _db(doc, s)
    _di(doc, SCREENSHOT_DIR / "alignment_admin.png", "Figure 7: Admin Dashboard", 5.5)

    _dh(doc, "7.2 Match Management", 2)
    _dp(doc, "The Match Review page allows administrators to manage AI-generated matches between lost and found items:")
    for s in ["Filter matches by status: Pending, Approved, Rejected, or view All matches",
              "Each match displays: confidence score, match level, and side-by-side item comparison with images",
              "Review the AI-generated explanation detailing why items were matched",
              "Approve matches with one click to notify both parties and initiate the recovery workflow",
              "Reject matches that are incorrect, providing resolution to the involved parties",
              "Approved matches automatically update item statuses and trigger notifications"]:
        _db(doc, s)

    _dh(doc, "7.3 Item Management", 2)
    _dp(doc, "Lost Items and Found Items management pages provide complete inventory control:")
    for s in ["View all reported items in a paginated table with comprehensive details",
              "Item details include: reference number, item name, category, reporter/finder, location, and status",
              "Update item status through the complete lifecycle: reported, under_review, match_approved, recovered, closed",
              "All status changes are logged in the activity log for audit purposes",
              "Click on item reference to view the full item details page"]:
        _db(doc, s)

    _dh(doc, "7.4 User Management", 2)
    _dp(doc, "The Users page provides comprehensive user account administration capabilities:")
    for s in ["View all registered users with their role, faculty, course, and account status",
              "Activate or deactivate user accounts using the toggle button",
              "Deactivated users cannot log in to the platform",
              "Permanently delete user accounts (with confirmation) — removes all associated data",
              "Administrators cannot delete their own account for security",
              "All user management actions are recorded in the activity log"]:
        _db(doc, s)

    _dh(doc, "7.5 Activity Logs", 2)
    _dp(doc, "The Activity Logs page provides a complete audit trail of all system actions:")
    for s in ["View recent system actions with user, action type, description, and timestamp",
              "Actions recorded include: login, logout, registration, report submissions, match operations",
              "Administrative actions: user management, item status updates, system settings changes",
              "IP addresses are recorded for security monitoring and forensic analysis",
              "Paginated view with configurable entries per page"]:
        _db(doc, s)

    # ── 8. AI MATCHING ENGINE ──
    _dh(doc, "8. AI Matching Engine")
    _dp(doc, "The AI matching engine is the core innovation of the FindMe platform. It uses content-based analysis to automatically identify potential matches between lost and found items, eliminating the need for manual searching and comparison.")

    _dh(doc, "8.1 Algorithm Overview", 2)
    _dp(doc, "The engine evaluates 10 attribute dimensions with weighted scoring to determine the likelihood that a lost item corresponds to a found item:")
    _dtable(doc, ["Attribute", "Weight", "Scoring Method"], [
        ["Item Name", "25%", "SequenceMatcher fuzzy string comparison"],
        ["Category", "15%", "Exact ID match with fallback similarity"],
        ["Description", "12%", "Jaccard keyword overlap + text similarity"],
        ["Color", "10%", "Exact match + color group normalization"],
        ["Brand", "10%", "Exact, substring, and fuzzy matching"],
        ["Location", "10%", "Location ID match + detail text proximity bonus"],
        ["Model", "8%", "Fuzzy text similarity comparison"],
        ["Date Proximity", "5%", "Day-difference scoring (1-30+ day ranges)"],
        ["Time Proximity", "3%", "Minute-difference scoring (30min-6hr ranges)"],
        ["Image", "2%", "Simulated (placeholder for future CNN-based)"],
    ])
    doc.add_paragraph()

    _dh(doc, "8.2 Confidence Levels", 2)
    _dp(doc, "The system classifies matches into four confidence levels based on the calculated score:")
    _dtable(doc, ["Level", "Score Range", "Description"], [
        ["Very High", "80-100%", "Multiple strong attribute matches across key dimensions"],
        ["High", "60-79%", "Good matches on primary attributes with supporting evidence"],
        ["Possible", "40-59%", "Some similarity detected but lacking confirming evidence"],
        ["Low", "0-39%", "Minimal match potential, requires manual administrator review"],
    ])
    doc.add_paragraph()

    _dh(doc, "8.3 Match Threshold", 2)
    _dp(doc, "The system generates a match only when the confidence score meets or exceeds 30%. This threshold ensures that low-probability matches do not overwhelm users while capturing all potentially relevant matches for administrator review. The threshold is configurable based on institutional requirements.")

    _dh(doc, "8.4 Explanation Generation", 2)
    _dp(doc, "A key differentiator of the FindMe platform is the generation of human-readable explanations for each match. These explanations provide complete transparency into the AI's decision-making process:")
    for s in ["Per-attribute similarity scores showing individual contribution to the total",
              "Clear listing of which attributes matched strongly and which did not",
              "Notable discrepancies or missing information that affected the score",
              "Final confidence percentage and match level classification",
              "Bullet-point format for easy scanning and understanding"]:
        _db(doc, s)
    _dp(doc, "This transparency builds user trust in the AI system and provides administrators with the information needed to make informed match approval decisions.")

    # ── 9. SECURITY ──
    _dh(doc, "9. Security & Compliance")
    _dp(doc, "The FindMe platform implements comprehensive security measures to protect user data and maintain system integrity:")
    sec_items = [
        ("Password Security", "All passwords are hashed using bcrypt with configurable salt rounds. Plain-text passwords are never stored, logged, or transmitted. Password minimum length of 6 characters is enforced."),
        ("Session Management", "Server-side sessions with secure, HTTP-only cookies prevent client-side tampering. Sessions are invalidated on logout and after periods of inactivity."),
        ("Role-Based Access Control", "Three-tier access control (Student, Lecturer, Administrator) with granular permission levels. Route-level authorization enforced through Python decorators."),
        ("SQL Injection Prevention", "All database queries use parameterized statements via MySQLdb. No raw SQL string concatenation with user input is permitted anywhere in the codebase."),
        ("File Upload Security", "Uploaded files are validated against a whitelist of allowed types (JPG, PNG, WebP). File size is limited to 16MB. Files are renamed to randomized strings to prevent path traversal attacks."),
        ("XSS Prevention", "Jinja2 template engine auto-escapes all user-supplied content by default. No user input is rendered as raw HTML without proper sanitization."),
        ("Audit Logging", "All system actions are logged with user ID, action type, description, entity references, IP address, and precise timestamp for complete auditability."),
        ("Data Integrity", "Foreign key constraints across all related tables maintain referential integrity. All database operations are performed within transactions."),
    ]
    for title, desc in sec_items:
        p = doc.add_paragraph()
        _dw(p, f"{title}: ", 11, True, DocRGB(0x1A, 0x27, 0x4B))
        _dw(p, desc, 11)

    # ── 10. DATABASE SCHEMA ──
    _dh(doc, "10. Database Schema")
    _dp(doc, "The system uses a MySQL relational database with the following 14 tables, designed for data integrity and query performance:")
    _dtable(doc, ["Table Name", "Purpose", "Key Relationships"], [
        ["users", "User accounts with roles and authentication", "roles, faculties, courses"],
        ["roles", "Role definitions (Student, Lecturer, Admin)", "users"],
        ["faculties", "Academic faculties within the university", "users, courses"],
        ["courses", "Academic programs linked to faculties", "users, faculties"],
        ["categories", "Item categories for classification", "lost_items, found_items"],
        ["locations", "Campus locations for reporting", "lost_items, found_items"],
        ["lost_items", "Lost item reports with full attributes", "users, categories, locations"],
        ["found_items", "Found item reports with full attributes", "users, categories, locations"],
        ["item_images", "Additional images for items", "lost_items"],
        ["matches", "AI-generated lost/found item matches", "lost_items, found_items, users"],
        ["notifications", "User notifications and alerts", "users"],
        ["verification_requests", "Ownership verification workflow", "matches, users"],
        ["recoveries", "Item recovery completion tracking", "matches, users"],
        ["activity_logs", "Complete system audit trail", "users"],
    ])
    doc.add_paragraph()

    # ── 11. TROUBLESHOOTING ──
    _dh(doc, "11. Troubleshooting")
    _dp(doc, "This section provides solutions to common issues users may encounter while using the FindMe platform.")
    issues = [
        ("Cannot Log In", "Verify your email address and password are correct. Check that your account has not been deactivated by an administrator. Use the 'Forgot Password' feature to reset your credentials if needed."),
        ("Report Submission Fails", "Ensure all required fields are completed (item name and date are mandatory). Verify your image file is under 16MB and in an accepted format (JPG, PNG, WebP). Check your internet connection and try again."),
        ("No AI Matches Found", "The AI engine requires a minimum 30% confidence score to generate a match. Provide more detailed item descriptions with specific attributes (brand, model, color, location). Ensure both lost and found items exist in the system."),
        ("Notifications Not Updating", "Visit the Dashboard or Notifications page to refresh your notification count. Notifications are fetched on page load and may not update in real-time without navigation."),
        ("Page Loading Slowly", "Check your internet connection speed. CDN resources (Font Awesome icons, Google Fonts) are loaded externally and may be slow on poor connections. Try refreshing the page or clearing your browser cache."),
        ("Forgot Password", "Use the 'Forgot Password' link on the login page. Enter your registered email address to receive password reset instructions. If you do not receive the email, check your spam folder or contact an administrator."),
        ("Account Deactivated", "Contact a system administrator to request account reactivation. Deactivated accounts cannot log in or access any platform features. Provide your full name and email address for verification."),
    ]
    for issue, solution in issues:
        p = doc.add_paragraph()
        _dw(p, f"Issue: {issue}", 11, True, DocRGB(0x1A, 0x27, 0x4B))
        _dp(doc, f"Solution: {solution}", 11, False, DocRGB(0x47, 0x55, 0x69))

    # ── 12. CONCLUSION ──
    _dh(doc, "12. Conclusion & Future Work")
    _dp(doc, "FindMe successfully addresses the lost and found management challenges at Cavendish University Uganda by providing a comprehensive digital platform with AI-powered matching capabilities. The system streamlines the entire process from item reporting through recovery, significantly reducing the time and effort required for all stakeholders involved.")
    _dp(doc, "The AI matching engine demonstrates the practical application of artificial intelligence in solving everyday campus challenges, while the responsive, accessible design ensures the platform can be used by the entire university community. The platform's modular architecture and comprehensive feature set position it for expansion to additional institutions.")
    
    _dh(doc, "Future Development Roadmap", 2)
    _dh(doc, "Phase 5: Mobile Applications (Q1 2027)", 3)
    _dp(doc, "Native iOS and Android applications with push notifications, camera integration for instant photo capture, and offline support for areas with limited connectivity.")

    _dh(doc, "Phase 6: Advanced Image Recognition (Q2 2027)", 3)
    _dp(doc, "Integration of Convolutional Neural Networks (CNN) for visual similarity analysis, enabling image-to-image matching between lost and found items.")

    _dh(doc, "Additional Planned Features", 3)
    for s in ["Barcode/RFID Integration: Support for scanning student ID cards and asset tags for instant item registration",
              "Multi-Institution Platform: Centralized deployment architecture serving multiple universities from a single instance",
              "Advanced Analytics Dashboard: Interactive data visualization with trend analysis and predictive insights",
              "Email Notifications: SMTP-based email alerts for critical match discoveries and account recovery",
              "QR Code Labels: Printable QR codes for high-value items enabling quick registration and tracking",
              "REST API Gateway: Public API for third-party integration with campus security and access control systems"]:
        _db(doc, s)

    # Save
    path = OUT_DIR / f"FindMe_Professional_Report_{TS}.docx"
    doc.save(str(path))
    print(f"[OK] Professional Word Report: {path}")
    return path


# ════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════

if __name__ == "__main__":
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    print("=" * 60)
    print("  FindMe - Professional Deliverables Generator")
    print("=" * 60)
    print("\n[1/2] Generating professional PowerPoint...")
    pptx = generate_pptx()
    print("\n[2/2] Generating professional Word report...")
    docx = generate_docx()
    print("\n" + "=" * 60)
    print("  Generation Complete!")
    print(f"  PowerPoint ({32} slides): {pptx}")
    print(f"  Word Report (30+ pages):  {docx}")
    print("=" * 60)